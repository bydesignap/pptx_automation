from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from typing import Dict, List, Tuple, Optional, Union
from pptx.shapes.base import BaseShape
from pptx.shapes.autoshape import Shape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.table import _Cell, Table
from io import BytesIO

class PresentationBuilder:
    """A class to build financial presentations with consistent styling"""
    
    def __init__(self, file_path=None):
        if file_path:
            self.prs = Presentation(file_path)
        else:
            self.prs = Presentation()
    
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        
        # Define common styles
        self.title_font = 'Arial'
        self.body_font = 'Arial'
        self.primary_color = RGBColor(0, 70, 127)  # Dark blue
        self.secondary_color = RGBColor(128, 128, 128)  # Gray
        
    def _apply_text_style(self, shape, font_size=12, font_name=None, bold=False, color=None):
        """Apply consistent text styling to a shape with type checking
        
        Args:
            shape: A shape object that may have a text frame
            font_size: Size of the font in points
            font_name: Name of the font to use
            bold: Whether the text should be bold
            color: RGB color for the text
            
        Returns:
            bool: True if styling was applied successfully, False otherwise
        """
        try:
            # Check if shape has text_frame attribute
            if not hasattr(shape, 'text_frame'):
                return False
                
            text_frame = shape.text_frame
            
            # Check if text_frame has paragraphs
            if not hasattr(text_frame, 'paragraphs') or len(text_frame.paragraphs) == 0:
                return False
                
            paragraph = text_frame.paragraphs[0]
            if not hasattr(paragraph, 'font'):
                return False
                
            # Apply text styling
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = font_name or self.body_font
            paragraph.font.bold = bold
            if color and hasattr(paragraph.font, 'color'):
                paragraph.font.color.rgb = color
                
            return True
        except Exception as e:
            print(f"Warning: Could not apply text style to shape: {str(e)}")
            return False
    
    def add_rounded_textbox(self, slide, title: str, content: str, left: float, top: float, width: float, height: float) -> None:
        """Create a text box with rounded corners and a table
        
        Args:
            title: The title for the first cell
            content: The content for the second cell
            left: The left position of the text box
            top: The top position of the text box
            width: The width of the text box
            height: The height of the text box
        """
        if not slide:
            raise ValueError("Slide must not be empty")
        if not title:
            raise ValueError("Title must not be empty")    
        if not left and not top and not width and not height:
            raise ValueError("Left, top, width or height must not be empty")    


        # Add rounded rectangle shape
        shape = slide.(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        
        # Set the corner radius to 2 points
        shape.adjustments[0] = 0.1  # 0.1 corresponds to 2 points
        
        # Add table inside the shape
        table = shape.text_frame.add_table(2, 1, Inches(left), Inches(top), Inches(width), Inches(height)).table
        
        # Set title cell
        cell = table.cell(0, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.primary_color
        self._apply_text_style(
            cell.text_frame,
            font_size=12,
            bold=True,
            color=RGBColor(255, 255, 255)
        )
        cell.text = title
        
        # Set content cell
        cell = table.cell(1, 0)
        self._apply_text_style(
            cell.text_frame,
            font_size=12,
            color=self.secondary_color
        )
        cell.text = content

    def create_custom_slide(self, title: str, logo_path: str) -> None:
        """Create a custom slide with a title and a logo
        
        Args:
            title: The slide title
            logo_path: Path to the logo file (SVG format)
            
        Raises:
            ValueError: If title is empty or logo_path is empty
        """
        if not title:
            raise ValueError("Title must not be empty")
        if not logo_path:
            raise ValueError("Logo path must not be empty")
        
        # Create a blank slide
        slide_layout = self.prs.slide_layouts[3]  # Using a blank layout
        self.prs.slide_width = Inches(13.33)  # Standard widescreen width
        self.prs.slide_height = Inches(7.5)
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.45),
            Inches(8), Inches(1)
        )

        try:
            title_box = slide.placeholders.title
        except AttributeError as e:
            title_box = slide.placeholders[0]
        
        self._apply_text_style(
            title_box,
            font_size=24,
            font_name=self.title_font,
            color=RGBColor(115, 139, 171)
        )
        title_box.text = title

        # Add logo
        try:
            logo = slide.shapes.add_picture(
                logo_path,
                Inches(12.3), Inches(0.14),
                Inches(.71), Inches(.71)
            )
        except Exception as e:
            print(f"Error adding logo: {str(e)}")

        return slide

    def create_title_slide(self, title: str, subtitle: str) -> None:
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_box = slide.shapes.title

        self._apply_text_style(
            title_box, 
            font_size=28, 
            font_name=self.title_font,
            bold=True,
            color=self.primary_color
        )
        title_box.text = title
        
        subtitle_box = slide.placeholders[0]

        self._apply_text_style(
            subtitle_box,
            font_size=18,
            color=self.secondary_color
        )

        subtitle_box.text = subtitle

        return slide
        
    def create_table_slide(self, title: str, headers: List[str], data: List[List[str]]) -> None:
        """Create a slide with a table
        
        Args:
            title: The slide title
            headers: List of column headers
            data: 2D list of table data
            
        Raises:
            ValueError: If title is empty, headers is empty, or data is empty
            ValueError: If data rows don't match header length
        """
        if not title:
            raise ValueError("Title must not be empty")
        if not headers:
            raise ValueError("Headers must not be empty")
        if not data:
            raise ValueError("Data must not be empty")
        
        # Verify data structure
        for row in data:
            if len(row) != len(headers):
                raise ValueError(f"Data row length {len(row)} does not match headers length {len(headers)}")
        """Create a slide with a table"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.75)
        )
        self._apply_text_style(
            title_box,
            font_size=20,
            bold=True,
            color=self.primary_color
        )
        title_box.text = title
        
        # Create table
        rows = len(data) + 1  # +1 for header
        cols = len(headers)
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), Inches(1.5),
            Inches(9), Inches(5)
        ).table
        
        # Set header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            self._apply_text_style(
                cell.text_frame,
                font_size=12,
                bold=True,
                color=self.primary_color
            )
            cell.text = header
            
        # Fill data
        for row_idx, row_data in enumerate(data, start=1):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                self._apply_text_style(
                    cell.text_frame,
                    font_size=11
                )
                cell.text = str(value)
                
    def save(self, filename: str) -> bool:
        """Save the presentation with error handling
        
        Args:
            filename: The filename to save the presentation to
            
        Returns:
            bool: True if save was successful, False otherwise
            
        Raises:
            ValueError: If filename is empty
        """
        if not filename:
            raise ValueError("Filename must not be empty")
            
        try:
            self.prs.save(filename)
            return True
        except Exception as e:
            print(f"Error saving presentation: {str(e)}")
            return False
        """Save the presentation"""
        self.prs.save(filename)

# Example usage
def create_presentation():
    builder = PresentationBuilder('template.pptx')
    
    # Title slide
    title_slide = builder.create_title_slide(
        "Third Quarter 2024\nEarnings Results Presentation",
        "October 15, 2024"
    )

    custom_slide = builder.create_custom_slide("Monitoring Dashboard", "/Users/alexparker/dev/python_test/GS.PNG")

    builder.add_rounded_textbox(custom_slide, "Total Revenue", "$2.5M", 1, 1, 4, 2)
    
    builder.save("earnings_presentation.pptx")

if __name__ == "__main__":
    create_presentation()